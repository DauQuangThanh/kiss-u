#!/usr/bin/env python3
"""Rebuild a .pptx from extended Markdown emitted by pptx_to_md.py.

Usage:
    python md_to_pptx.py input.md output.pptx [--template original.pptx] [--meta deck.meta.json]

Round-trip works best when the original .pptx is supplied via --template so the
master slides, layouts, theme, and chart parts carry over. If no template is
specified, a blank default deck is used and layouts are matched by name as best
as possible.
"""
from __future__ import annotations

import sys
from pathlib import Path

# Run the venv bootstrap BEFORE any third-party imports. If ./.venv exists in
# the user's cwd, this re-execs the script under that interpreter; if not, it
# prints setup instructions unless --system-python was passed.
if __name__ == "__main__":
    sys.path.insert(0, str(Path(__file__).parent))
    from _venv_bootstrap import ensure_env  # noqa: E402
    ensure_env(__file__, skill_name="kiss-pptx-markdown")

import argparse  # noqa: E402
import copy  # noqa: E402
import json  # noqa: E402
import os  # noqa: E402
import re  # noqa: E402


def _pip_hint(packages: str) -> str:
    """Return a pip install hint that works on Windows, macOS, and Linux."""
    if sys.platform.startswith("win"):
        return f"pip install {packages}"
    return f"pip install {packages} --break-system-packages"


try:
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.exit(f"Install python-pptx: {_pip_hint('python-pptx pyyaml')}")

try:
    import yaml
except ImportError:
    sys.exit(f"Install pyyaml: {_pip_hint('pyyaml')}")


def _open_template_as_blank(tpl_path: Path):
    """Open a template .pptx and strip all existing slides cleanly.

    python-pptx retains slide parts in the package even after we remove their
    references from sldIdLst, which produces duplicate-name warnings and can
    confuse PowerPoint. We strip them directly from the package.
    """
    prs = Presentation(str(tpl_path))
    slide_parts = list(prs.slides._sldIdLst)
    # Remove each slide from the presentation's rId list
    for sldId in slide_parts:
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.slides._sldIdLst.remove(sldId)
        if rId and rId in prs.part.rels:
            # Drop the slide part from the package
            slide_part = prs.part.rels[rId].target_part
            prs.part.drop_rel(rId)
            try:
                prs.part.package._parts.pop(slide_part.partname, None)
            except Exception:
                pass
    return prs


SLIDE_START_RE = re.compile(r"<!--\s*slide:start\s+(?P<attrs>[^>]+?)-->")
SLIDE_END_RE = re.compile(r"<!--\s*slide:end\s*-->")
NOTES_START_RE = re.compile(r"<!--\s*notes:start\s*-->")
NOTES_END_RE = re.compile(r"<!--\s*notes:end\s*-->")
IMAGE_RE = re.compile(r"!\[(?P<alt>[^\]]*)\]\((?P<src>[^)]+)\)(?:\{(?P<attrs>[^}]*)\})?")
SHAPE_REF_RE = re.compile(r"<!--\s*shape\s+id=(?P<id>\S+)\s+kind=(?P<kind>\S+)[^>]*-->")


def parse_attr_block(attr_str: str) -> dict:
    """Parse a Pandoc-style attribute block or HTML-comment attribute list.

    Supports `#id`, `.class`, `key=value`, and `key="value with spaces"`.
    """
    out: dict = {"classes": []}
    i = 0
    n = len(attr_str)
    while i < n:
        while i < n and attr_str[i].isspace():
            i += 1
        if i >= n:
            break
        ch = attr_str[i]
        if ch == "#":
            j = i + 1
            while j < n and not attr_str[j].isspace():
                j += 1
            out["id"] = attr_str[i + 1 : j]
            i = j
        elif ch == ".":
            j = i + 1
            while j < n and not attr_str[j].isspace():
                j += 1
            out["classes"].append(attr_str[i + 1 : j])
            i = j
        else:
            # key=value, value may be quoted
            j = i
            while j < n and attr_str[j] not in "= \t":
                j += 1
            key = attr_str[i:j]
            if j >= n or attr_str[j] != "=":
                i = j
                continue
            j += 1  # skip '='
            if j < n and attr_str[j] == '"':
                k = attr_str.find('"', j + 1)
                if k == -1:
                    k = n
                out[key] = attr_str[j + 1 : k]
                i = k + 1
            else:
                k = j
                while k < n and not attr_str[k].isspace():
                    k += 1
                out[key] = attr_str[j:k]
                i = k
    return out


def unit_to_emu(value: str) -> int:
    m = re.match(r"^([0-9]*\.?[0-9]+)(in|cm|pt)?$", value)
    if not m:
        return 0
    num = float(m.group(1))
    unit = m.group(2) or "in"
    if unit == "in":
        return int(num * 914400)
    if unit == "cm":
        return int(num * 360000)
    if unit == "pt":
        return int(num * 12700)
    return int(num)


def parse_frontmatter(text: str) -> tuple[dict, str]:
    if not text.startswith("---"):
        return {}, text
    end = text.find("\n---", 3)
    if end == -1:
        return {}, text
    fm_text = text[3:end].strip()
    body = text[end + 4 :].lstrip("\n")
    return yaml.safe_load(fm_text) or {}, body


def split_slides(body: str) -> list[dict]:
    slides: list[dict] = []
    pos = 0
    while True:
        m = SLIDE_START_RE.search(body, pos)
        if not m:
            break
        end_m = SLIDE_END_RE.search(body, m.end())
        if not end_m:
            raise ValueError("Unclosed <!-- slide:start --> block")
        attrs = parse_attr_block(m.group("attrs"))
        content = body[m.end() : end_m.start()]
        notes = ""
        ns = NOTES_START_RE.search(content)
        if ns:
            ne = NOTES_END_RE.search(content, ns.end())
            if ne:
                notes = content[ns.end() : ne.start()].strip()
                content = content[: ns.start()] + content[ne.end() :]
        slides.append(
            {
                "layout": attrs.get("layout", ""),
                "id": attrs.get("id", str(len(slides) + 1)),
                "background": attrs.get("background"),
                "hidden": attrs.get("hidden") == "true",
                "body": content.strip(),
                "notes": notes,
            }
        )
        pos = end_m.end()
    return slides


def find_layout(prs: Presentation, name: str):
    name = (name or "").strip()
    for layout in prs.slide_layouts:
        if (layout.name or "") == name:
            return layout
    # Fallback: first layout with a title placeholder.
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.type in (13, 15):
                return layout
    return prs.slide_layouts[0]


def extract_title_subtitle(body: str) -> tuple[str, str, str]:
    """Return (title, subtitle, remaining_body)."""
    title = ""
    subtitle = ""
    lines = body.splitlines()
    remaining: list[str] = []
    consumed_title = False
    consumed_subtitle = False
    for line in lines:
        if not consumed_title and line.startswith("# "):
            title = line[2:].strip()
            consumed_title = True
            continue
        if consumed_title and not consumed_subtitle and line.startswith("## "):
            subtitle = line[3:].strip()
            consumed_subtitle = True
            continue
        remaining.append(line)
    return title, subtitle, "\n".join(remaining).strip()


def set_placeholder_text(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    if not text:
        return
    lines = text.splitlines() or [""]
    p = tf.paragraphs[0]
    p.text = lines[0]
    for line in lines[1:]:
        np = tf.add_paragraph()
        np.text = line


def add_body_content(slide, body: str, meta_slide: dict | None, md_dir: Path) -> None:
    """Write list/paragraph/table/image content into a body placeholder (best effort)."""
    # Find a body placeholder (non-title, non-subtitle).
    body_ph = None
    for ph in slide.placeholders:
        pt = ph.placeholder_format.type
        if pt in (1, 3, 4, 5):
            continue
        body_ph = ph
        break

    # Parse body into blocks (paragraph, list item, image, table, shape-ref).
    blocks = _blocks_from_markdown(body)
    text_blocks = [b for b in blocks if b["type"] in ("paragraph", "bullet", "number")]
    media_blocks = [b for b in blocks if b["type"] in ("image", "table", "shape_ref", "code")]

    # Write text blocks into body placeholder.
    if body_ph is not None and text_blocks:
        tf = body_ph.text_frame
        tf.clear()
        first = True
        for b in text_blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = b["text"]
            if b["type"] == "bullet":
                p.level = b.get("level", 0)
                _ensure_bullet(p, numbered=False)
            elif b["type"] == "number":
                p.level = b.get("level", 0)
                _ensure_bullet(p, numbered=True)
            else:
                _ensure_bullet(p, numbered=False, disable=True)
    elif text_blocks:
        # No body placeholder; create a text box.
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        first = True
        for b in text_blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = b["text"]

    # Handle media blocks.
    for b in media_blocks:
        if b["type"] == "image":
            _add_image(slide, b, meta_slide, md_dir)
        elif b["type"] == "table":
            _add_table(slide, b["rows"])
        elif b["type"] == "code":
            _add_code_block(slide, b["text"])
        elif b["type"] == "shape_ref":
            # Charts, smartart etc. recreated from template chart parts are
            # already in the slide because we started from the template.
            # Nothing to do here unless sidecar has a different shape to inject.
            pass


def _ensure_bullet(paragraph, numbered: bool, disable: bool = False) -> None:
    """Force bullet/number formatting on a paragraph at the DrawingML level.

    Body placeholders inherit bullets from their layout; this function makes
    the intent explicit so it round-trips even when the layout's default
    doesn't match what the Markdown says.
    """
    from lxml import etree

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_el = paragraph._p
    pPr = p_el.find(f"{{{a_ns}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(p_el, f"{{{a_ns}}}pPr")
        p_el.insert(0, pPr)
    # Remove any existing bullet/number markers.
    for tag in ("buChar", "buAutoNum", "buNone"):
        existing = pPr.find(f"{{{a_ns}}}{tag}")
        if existing is not None:
            pPr.remove(existing)
    if disable:
        etree.SubElement(pPr, f"{{{a_ns}}}buNone")
        return
    if numbered:
        el = etree.SubElement(pPr, f"{{{a_ns}}}buAutoNum")
        el.set("type", "arabicPeriod")
    else:
        el = etree.SubElement(pPr, f"{{{a_ns}}}buChar")
        el.set("char", "\u2022")


def _blocks_from_markdown(body: str) -> list[dict]:
    blocks: list[dict] = []
    lines = body.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        if not line.strip():
            i += 1
            continue
        m_img = IMAGE_RE.match(line.strip())
        if m_img:
            attrs = parse_attr_block(m_img.group("attrs") or "")
            blocks.append({
                "type": "image",
                "alt": m_img.group("alt"),
                "src": m_img.group("src"),
                "attrs": attrs,
            })
            i += 1
            continue
        m_ref = SHAPE_REF_RE.match(line.strip())
        if m_ref:
            blocks.append({"type": "shape_ref", "id": m_ref.group("id"), "kind": m_ref.group("kind")})
            i += 1
            continue
        if line.startswith("```"):
            # code block
            fence = line
            i += 1
            buf: list[str] = []
            while i < len(lines) and not lines[i].startswith("```"):
                buf.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1  # closing fence
            blocks.append({"type": "code", "text": "\n".join(buf)})
            continue
        if line.lstrip().startswith("|") and i + 1 < len(lines) and re.match(r"^\|?\s*[:\- |]+\|?\s*$", lines[i + 1]):
            # GFM table
            rows = [line]
            i += 1
            rows.append(lines[i])  # separator
            i += 1
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                rows.append(lines[i])
                i += 1
            blocks.append({"type": "table", "rows": _parse_gfm_table(rows)})
            continue
        # list item?
        m = re.match(r"^(\s*)([-*])\s+(.*)$", line)
        if m:
            level = len(m.group(1)) // 2
            blocks.append({"type": "bullet", "level": level, "text": m.group(3)})
            i += 1
            continue
        m = re.match(r"^(\s*)(\d+)\.\s+(.*)$", line)
        if m:
            level = len(m.group(1)) // 2
            blocks.append({"type": "number", "level": level, "text": m.group(3)})
            i += 1
            continue
        blocks.append({"type": "paragraph", "text": line.rstrip()})
        i += 1
    return blocks


def _parse_gfm_table(rows: list[str]) -> list[list[str]]:
    def split_row(r: str) -> list[str]:
        r = r.strip()
        if r.startswith("|"):
            r = r[1:]
        if r.endswith("|"):
            r = r[:-1]
        return [c.strip().replace("<br>", "\n") for c in r.split("|")]

    header = split_row(rows[0])
    data_rows = [split_row(r) for r in rows[2:]]
    return [header] + data_rows


def _add_image(slide, block: dict, meta_slide: dict | None, md_dir: Path) -> None:
    src = block["src"]
    path = (md_dir / src).resolve()
    if not path.exists():
        print(f"warning: image not found: {path}", file=sys.stderr)
        return
    attrs = block.get("attrs", {})
    sid = attrs.get("id")
    left = top = width = height = None
    if meta_slide and sid and sid in meta_slide.get("shapes", {}):
        s = meta_slide["shapes"][sid]
        left = s.get("left_emu")
        top = s.get("top_emu")
        width = s.get("cx_emu")
        height = s.get("cy_emu")
    if "x" in attrs:
        left = unit_to_emu(attrs["x"])
    if "y" in attrs:
        top = unit_to_emu(attrs["y"])
    if "width" in attrs:
        width = unit_to_emu(attrs["width"])
    if "height" in attrs:
        height = unit_to_emu(attrs["height"])
    if any(v is None for v in (left, top, width, height)):
        # fall back to sensible defaults
        left = left or Inches(1)
        top = top or Inches(1.5)
        width = width or Inches(6)
        height = height or Inches(3.5)
    slide.shapes.add_picture(str(path), Emu(left), Emu(top), width=Emu(width), height=Emu(height))


def _add_table(slide, rows: list[list[str]]) -> None:
    if not rows:
        return
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(0.4 * nrows + 0.2)
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = table_shape.table
    for r_idx, r in enumerate(rows):
        for c_idx in range(ncols):
            cell = table.cell(r_idx, c_idx)
            cell.text = r[c_idx] if c_idx < len(r) else ""


def _add_code_block(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.name = "Consolas"
        run.font.size = Pt(12)


def set_notes(slide, notes: str) -> None:
    if not notes:
        notes = "(No speaker notes. Add them before presenting.)"
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    first = True
    for line in notes.splitlines() or [""]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line


def rebuild(md_path: Path, out_path: Path, template: Path | None, meta_path: Path | None) -> None:
    text = md_path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)

    # Resolve template.
    tpl_ref = template or (fm.get("template") and md_path.parent / fm["template"])
    if tpl_ref and Path(tpl_ref).exists():
        prs = _open_template_as_blank(Path(tpl_ref))
    else:
        prs = Presentation()
        ss = fm.get("slide_size") or {}
        if ss:
            prs.slide_width = Emu(int(unit_to_emu(f"{ss.get('width', 13.333)}{ss.get('unit', 'in')}")))
            prs.slide_height = Emu(int(unit_to_emu(f"{ss.get('height', 7.5)}{ss.get('unit', 'in')}")))

    # Load sidecar.
    meta: dict = {}
    if meta_path is None:
        guess = md_path.with_suffix("").with_suffix(".meta.json")
        if fm.get("pptx_meta"):
            guess = md_path.parent / fm["pptx_meta"]
        if guess.exists():
            meta_path = guess
    if meta_path and Path(meta_path).exists():
        meta = json.loads(Path(meta_path).read_text(encoding="utf-8"))

    slides = split_slides(body)
    md_dir = md_path.parent

    for s in slides:
        layout = find_layout(prs, s["layout"])
        slide = prs.slides.add_slide(layout)
        title, subtitle, rest = extract_title_subtitle(s["body"])
        # Title + subtitle placeholders
        for ph in slide.placeholders:
            pt = ph.placeholder_format.type
            if pt in (1, 3, 5) and title:
                set_placeholder_text(ph, title)
            elif pt == 4 and subtitle:
                set_placeholder_text(ph, subtitle)
        meta_slide = meta.get("slides", {}).get(s["id"])
        add_body_content(slide, rest, meta_slide, md_dir)
        set_notes(slide, s["notes"])

    # Core props.
    if fm.get("title"):
        prs.core_properties.title = fm["title"]
    if fm.get("author"):
        prs.core_properties.author = fm["author"]

    prs.save(str(out_path))
    print(f"Wrote {out_path}")


def main() -> None:
    ap = argparse.ArgumentParser(description="Rebuild a .pptx from Markdown.")
    ap.add_argument("input", type=Path)
    ap.add_argument("output", type=Path)
    ap.add_argument("--template", type=Path, default=None)
    ap.add_argument("--meta", type=Path, default=None)
    args = ap.parse_args()
    rebuild(args.input, args.output, args.template, args.meta)


if __name__ == "__main__":
    main()
