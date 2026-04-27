#!/usr/bin/env python3
"""Export a .pptx file to extended Markdown + assets + meta.json sidecar.

Usage:
    python pptx_to_md.py input.pptx output.md [--assets-dir DIR] [--no-meta]
    python pptx_to_md.py input.pptx --list-layouts

See ../references/format-spec.md for the Markdown dialect this emits.
"""
from __future__ import annotations

import sys
from pathlib import Path

# Run the venv bootstrap BEFORE any third-party imports. If ./.venv exists in
# the user's cwd, this re-execs the script under that interpreter; if not, it
# prints setup instructions unless --system-python was passed.
if __name__ == "__main__":
    sys.path.insert(0, str(Path(__file__).parent))
    from venv_bootstrap import ensure_env  # noqa: E402
    ensure_env(__file__, skill_name="kiss-pptx-markdown")

import argparse  # noqa: E402
import json  # noqa: E402
import os  # noqa: E402
import re  # noqa: E402


def _pip_hint(packages: str) -> str:
    """Return a pip install hint that works on Windows, macOS, and Linux."""
    if sys.platform.startswith("win"):
        return f"pip install {packages}"
    return f"pip install {packages} --break-system-packages"


try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        from pptx.enum.text import PP_PLACEHOLDER  # older versions
except ImportError:
    sys.exit(f"Install python-pptx: {_pip_hint('python-pptx pyyaml')}")


TITLE_TYPES = {1, 3, 5}   # TITLE, CENTER_TITLE, VERTICAL_TITLE
SUBTITLE_TYPES = {4}       # SUBTITLE
BODY_TYPES = {2, 6, 7}     # BODY, VERTICAL_BODY, OBJECT

try:
    import yaml
except ImportError:
    sys.exit(f"Install pyyaml: {_pip_hint('pyyaml')}")


EMU_PER_INCH = 914400


def emu_to_inches(emu: int) -> float:
    return round(emu / EMU_PER_INCH, 4)


def escape_md(text: str) -> str:
    if not text:
        return ""
    # Escape characters that would otherwise change Markdown meaning.
    return re.sub(r"([\\`*_{}\[\]()#+\-!<>])", r"\\\1", text)


def text_frame_to_markdown_list(tf, force_bullets: bool = False) -> str:
    """Body-placeholder version: treats paragraphs as bullets by default.

    Body placeholders in PowerPoint default to a bulleted style that's defined
    on the layout/master via `buAutoNum` or `buChar` but is not always present
    on the per-paragraph `<a:pPr>`. We detect this by inspecting the text
    frame's parent shape placeholder type; if the caller says so, treat every
    non-empty paragraph as a bullet at its indent level.
    """
    if tf is None:
        return ""
    lines: list[str] = []
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs) or para.text
        text = text.rstrip()
        if not text:
            lines.append("")
            continue
        level = para.level or 0
        indent = "  " * level
        explicit_bullet, explicit_number = _paragraph_marker_type(para)
        if explicit_number:
            lines.append(f"{indent}1. {text}")
        elif explicit_bullet or force_bullets:
            lines.append(f"{indent}- {text}")
        else:
            lines.append(f"{indent}{text}")
    out = "\n".join(lines).rstrip()
    return out + "\n" if out else ""


def _paragraph_marker_type(para) -> tuple[bool, bool]:
    pPr = para._pPr
    if pPr is None:
        return False, False
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    if pPr.find(f"{a_ns}buNone") is not None:
        return False, False
    if pPr.find(f"{a_ns}buAutoNum") is not None:
        return False, True
    if pPr.find(f"{a_ns}buChar") is not None:
        return True, False
    return False, False


def text_frame_to_markdown(tf) -> str:
    """Convert a python-pptx text frame to a Markdown-ish string.

    Preserves paragraphs, bullets (as `-`), and numbered lists (as `1.`).
    Ignores font/color styling: those come from the master/layout.
    """
    if tf is None:
        return ""
    lines: list[str] = []
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs) or para.text
        text = text.rstrip()
        if not text:
            lines.append("")
            continue
        level = para.level or 0
        indent = "  " * level
        pPr = para._pPr
        bullet = False
        numbered = False
        if pPr is not None:
            if pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None:
                bullet = True
            elif pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum") is not None:
                numbered = True
        if bullet:
            lines.append(f"{indent}- {text}")
        elif numbered:
            lines.append(f"{indent}1. {text}")
        else:
            lines.append(f"{indent}{text}")
    return "\n".join(lines).rstrip() + "\n"


def table_to_markdown(table) -> str:
    """Convert a python-pptx table to a GFM table."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            txt = " ".join(p.text for p in cell.text_frame.paragraphs).strip()
            txt = txt.replace("|", "\\|").replace("\n", "<br>")
            cells.append(txt)
        rows.append(cells)
    if not rows:
        return ""
    ncols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < ncols:
            r.append("")
    out = ["| " + " | ".join(rows[0]) + " |"]
    out.append("|" + "|".join([" --- "] * ncols) + "|")
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out) + "\n"


def export_image(shape, assets_dir: Path, shape_id: str) -> tuple[str, str]:
    """Save the image to assets_dir, return (relative_path, ext)."""
    image = shape.image
    ext = image.ext or "png"
    filename = f"{shape_id}.{ext}"
    out_path = assets_dir / filename
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "wb") as f:
        f.write(image.blob)
    return str(out_path.relative_to(assets_dir.parent)), ext


def shape_position_attrs(shape) -> dict:
    return {
        "left_emu": int(shape.left) if shape.left is not None else 0,
        "top_emu": int(shape.top) if shape.top is not None else 0,
        "cx_emu": int(shape.width) if shape.width is not None else 0,
        "cy_emu": int(shape.height) if shape.height is not None else 0,
        "rotation": float(shape.rotation) if hasattr(shape, "rotation") else 0.0,
    }


def attrs_to_inline(shape, shape_id: str, klass: str = "") -> str:
    """Pandoc-style attribute block: {#id .class width=... height=... x=... y=...}"""
    pos = shape_position_attrs(shape)
    parts = [f"#{shape_id}"]
    if klass:
        parts.append(f".{klass}")
    parts.append(f"width={emu_to_inches(pos['cx_emu'])}in")
    parts.append(f"height={emu_to_inches(pos['cy_emu'])}in")
    parts.append(f"x={emu_to_inches(pos['left_emu'])}in")
    parts.append(f"y={emu_to_inches(pos['top_emu'])}in")
    return "{" + " ".join(parts) + "}"


def iter_shapes_with_ids(slide):
    """Yield (shape_id, shape) in z-order. Uses shape.shape_id when stable."""
    for idx, shape in enumerate(slide.shapes, start=1):
        sid = f"shape-{shape.shape_id}" if getattr(shape, "shape_id", None) else f"shape-{idx}"
        yield sid, shape


def slide_to_markdown(slide, slide_idx: int, assets_dir: Path, meta: dict) -> str:
    layout_name = slide.slide_layout.name or "Blank"
    parts: list[str] = []
    parts.append(
        f'<!-- slide:start layout="{layout_name}" id="{slide_idx}" -->'
    )

    title_text = ""
    subtitle_text = ""
    body_blocks: list[str] = []
    shape_meta: dict = {}

    for sid, shape in iter_shapes_with_ids(slide):
        stype = shape.shape_type
        shape_meta[sid] = {
            "kind": "unknown",
            **shape_position_attrs(shape),
            "name": shape.name,
        }

        ph_type_val = None
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            ph_type_val = int(ph_type) if ph_type is not None else None

        if ph_type_val in TITLE_TYPES and not title_text and shape.has_text_frame:
            title_text = shape.text_frame.text.strip()
            shape_meta[sid]["kind"] = "placeholder_title"
            continue
        if ph_type_val in SUBTITLE_TYPES and not subtitle_text and shape.has_text_frame:
            subtitle_text = shape.text_frame.text.strip()
            shape_meta[sid]["kind"] = "placeholder_subtitle"
            continue

        if stype == MSO_SHAPE_TYPE.PICTURE or (_has_image(shape)):
            try:
                rel, ext = export_image(shape, assets_dir, sid)
                alt = shape.name or "image"
                attrs = attrs_to_inline(shape, sid, klass="picture")
                body_blocks.append(f"![{alt}]({rel}){attrs}")
                shape_meta[sid]["kind"] = "picture"
                shape_meta[sid]["asset"] = rel
                continue
            except Exception:
                pass

        if stype == MSO_SHAPE_TYPE.TABLE and shape.has_table:
            body_blocks.append(table_to_markdown(shape.table))
            shape_meta[sid]["kind"] = "table"
            continue

        if stype == MSO_SHAPE_TYPE.CHART:
            body_blocks.append(
                f"<!-- shape id={sid} kind=chart ref=meta.json#slides.{slide_idx}.shapes.{sid} -->"
            )
            shape_meta[sid]["kind"] = "chart"
            continue

        if shape.has_text_frame:
            # A body/content placeholder with bullets — convert to Markdown list.
            is_body_ph = ph_type_val in BODY_TYPES
            text = text_frame_to_markdown_list(shape.text_frame, force_bullets=is_body_ph)
            if text.strip():
                body_blocks.append(text.rstrip())
            shape_meta[sid]["kind"] = "placeholder_body" if is_body_ph else "text"
            continue

        body_blocks.append(
            f"<!-- shape id={sid} kind={_shape_kind_name(stype)} ref=meta.json#slides.{slide_idx}.shapes.{sid} -->"
        )
        shape_meta[sid]["kind"] = _shape_kind_name(stype)

    # Emit title/subtitle first, then body.
    if title_text:
        parts.append(f"\n# {title_text}\n")
    if subtitle_text:
        parts.append(f"## {subtitle_text}\n")
    if body_blocks:
        parts.append("\n" + "\n\n".join(b for b in body_blocks if b) + "\n")

    # Speaker notes
    notes_text = ""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
    if not notes_text:
        notes_text = "(No speaker notes in source. Add notes before presenting.)"
    parts.append("\n<!-- notes:start -->")
    parts.append(notes_text)
    parts.append("<!-- notes:end -->")
    parts.append("<!-- slide:end -->\n")

    # Record slide meta
    meta.setdefault("slides", {})[str(slide_idx)] = {
        "layout_name": layout_name,
        "hidden": bool(getattr(slide, "element", None) is not None
                       and slide.element.get("show") == "0"),
        "shapes": shape_meta,
    }
    return "\n".join(parts)


def _has_image(shape) -> bool:
    try:
        _ = shape.image
        return True
    except Exception:
        return False


def _shape_kind_name(stype) -> str:
    if stype is None:
        return "unknown"
    try:
        return str(stype).split(".")[-1].lower()
    except Exception:
        return "unknown"


def build_frontmatter(prs: Presentation, meta_path: str | None, template_path: str | None) -> str:
    cp = prs.core_properties
    fm = {
        "title": cp.title or "",
        "author": cp.author or "",
        "created": cp.created.isoformat() if cp.created else "",
        "modified": cp.modified.isoformat() if cp.modified else "",
        "slide_size": {
            "width": emu_to_inches(prs.slide_width),
            "height": emu_to_inches(prs.slide_height),
            "unit": "in",
        },
    }
    if template_path:
        fm["template"] = template_path
    if meta_path:
        fm["pptx_meta"] = meta_path
    return "---\n" + yaml.safe_dump(fm, sort_keys=True, default_flow_style=False).rstrip() + "\n---\n"


def list_layouts(prs: Presentation) -> None:
    for layout in prs.slide_layouts:
        print(layout.name)


def convert(src: Path, dest: Path, assets_dir: Path | None, write_meta: bool) -> None:
    prs = Presentation(str(src))
    basename = dest.stem
    if assets_dir is None:
        assets_dir = dest.parent / f"{basename}.assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    meta_path = dest.parent / f"{basename}.meta.json"

    meta = {
        "slide_size_emu": {"cx": prs.slide_width, "cy": prs.slide_height},
        "template_path": str(src.name),
        "core_properties": {
            "title": prs.core_properties.title or "",
            "author": prs.core_properties.author or "",
            "created": prs.core_properties.created.isoformat() if prs.core_properties.created else "",
            "modified": prs.core_properties.modified.isoformat() if prs.core_properties.modified else "",
        },
        "slides": {},
    }

    body_parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        body_parts.append(slide_to_markdown(slide, idx, assets_dir, meta))

    fm = build_frontmatter(
        prs,
        meta_path=str(meta_path.name) if write_meta else None,
        template_path=str(src.name),
    )
    body = "\n".join(body_parts)
    dest.write_text(fm + "\n" + body, encoding="utf-8")

    if write_meta:
        meta_path.write_text(json.dumps(meta, indent=2, sort_keys=True), encoding="utf-8")

    print(f"Wrote {dest}")
    print(f"  assets -> {assets_dir}")
    if write_meta:
        print(f"  meta   -> {meta_path}")


def main() -> None:
    ap = argparse.ArgumentParser(description="Convert .pptx to extended Markdown.")
    ap.add_argument("input", type=Path, help="Input .pptx file")
    ap.add_argument("output", type=Path, nargs="?", help="Output .md file")
    ap.add_argument("--assets-dir", type=Path, default=None)
    ap.add_argument("--no-meta", action="store_true")
    ap.add_argument("--list-layouts", action="store_true")
    args = ap.parse_args()

    if args.list_layouts:
        list_layouts(Presentation(str(args.input)))
        return

    if args.output is None:
        ap.error("output is required unless --list-layouts is used")

    convert(args.input, args.output, args.assets_dir, write_meta=not args.no_meta)


if __name__ == "__main__":
    main()
